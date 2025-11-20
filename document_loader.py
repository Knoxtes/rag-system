# document_loader.py - Load and process all file types with Shared Drive support and OCR

import io
import logging
from typing import Optional
from googleapiclient.http import MediaIoBaseDownload
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd
from config import (
    CHUNK_SIZE, 
    CHUNK_OVERLAP, 
    USE_PARENT_DOCUMENT_RETRIEVAL, 
    PARENT_CHUNK_SIZE,
    # OCR Settings
    OCR_BACKEND,
    OCR_CONFIDENCE_THRESHOLD,
    OCR_LANGUAGES,
    OCR_ENABLED,
    # Document AI Settings
    DOCUMENTAI_PROJECT_ID,
    DOCUMENTAI_LOCATION,
    DOCUMENTAI_PROCESSOR_ID
)
# --- OPTIMIZATION: Use Langchain for robust, standard chunning ---
from langchain_text_splitters import RecursiveCharacterTextSplitter
import re

# OCR Integration (optional)
try:
    from ocr_service import OCRServiceFactory, is_image_file, get_image_format_from_mime
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False
    def is_image_file(mime_type, filename): return False
    def get_image_format_from_mime(mime_type): return None
    class OCRServiceFactory:
        @staticmethod
        def create_service(*args, **kwargs): return None

# Text Clarification Integration (optional)
try:
    from text_clarification import get_clarification_service
    TEXT_CLARIFICATION_AVAILABLE = True
except ImportError:
    TEXT_CLARIFICATION_AVAILABLE = False
    def get_clarification_service(): return None

# Text Quality Filter Integration (optional)
try:
    from text_quality_filter import get_quality_filter
    TEXT_QUALITY_AVAILABLE = True
except ImportError:
    TEXT_QUALITY_AVAILABLE = False
    def get_quality_filter():
        class DummyFilter:
            def assess_text_quality(self, text, filename, source_type):
                return {'should_include': True, 'quality_score': 1.0, 'reason': 'Quality filter not available'}
        return DummyFilter()

# Setup logging
logger = logging.getLogger(__name__)


class GoogleDriveLoader:
    """Download and process files from Google Drive and Shared Drives"""
    
    def __init__(self, service):
        self.service = service
        self.ocr_service = None
        self._initialize_ocr()
    
    def _initialize_ocr(self):
        """Initialize OCR service if enabled"""
        if OCR_ENABLED:
            try:
                # Use Document AI if specified
                if OCR_BACKEND == "documentai":
                    from documentai_ocr import create_documentai_service
                    self.ocr_service = create_documentai_service(
                        project_id=DOCUMENTAI_PROJECT_ID,
                        location=DOCUMENTAI_LOCATION,
                        processor_id=DOCUMENTAI_PROCESSOR_ID,
                        confidence_threshold=OCR_CONFIDENCE_THRESHOLD
                    )
                    logger.info(f"Document AI OCR service initialized (project: {DOCUMENTAI_PROJECT_ID})")
                else:
                    # Fallback to legacy OCR service factory
                    self.ocr_service = OCRServiceFactory.create_service(
                        backend=OCR_BACKEND,
                        confidence_threshold=OCR_CONFIDENCE_THRESHOLD
                    )
                    logger.info(f"OCR service initialized: {OCR_BACKEND}")
            except Exception as e:
                logger.warning(f"Failed to initialize OCR service: {e}")
                self.ocr_service = None
    
    # --- REMOVED: list_files() method was unused (dead code) ---
    
    def download_file(self, file_id):
        """Download file content"""
        try:
            request = self.service.files().get_media(fileId=file_id, supportsAllDrives=True)
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            file_content.seek(0)
            return file_content
        except Exception as e:
            print(f"  Error downloading: {e}")
            return None
    
    def export_google_doc(self, file_id):
        """Export Google Docs as text"""
        try:
            request = self.service.files().export_media(
                fileId=file_id, 
                mimeType='text/plain'
            )
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            return file_content.getvalue().decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"  Error exporting: {e}")
            return None
    
    def export_google_slides(self, file_id):
        """Export Google Slides as text"""
        # Slides export to text/plain is surprisingly effective
        return self.export_google_doc(file_id)
    
    def export_google_sheets(self, file_id):
        """Export Google Sheets as CSV"""
        try:
            request = self.service.files().export_media(
                fileId=file_id, 
                mimeType='text/csv'
            )
            file_content = io.BytesIO()
            downloader = MediaIoBaseDownload(file_content, request)
            
            done = False
            while not done:
                status, done = downloader.next_chunk()
            
            return file_content.getvalue().decode('utf-8', errors='ignore')
        except Exception as e:
            print(f"  Error exporting: {e}")
            return None

# --- Text Extraction Functions (Largely Unchanged) ---

def extract_text_from_pptx(file_content):
    """Extract text from PowerPoint"""
    try:
        prs = Presentation(file_content)
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"[Slide {slide_num}]"]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Only add slide if it has text content
            if len(slide_text) > 1:
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"  Error extracting PPTX: {e}")
        return ""


def extract_text_from_xlsx(file_content):
    """Extract text from Excel - returns complete file without chunking marker"""
    try:
        df_dict = pd.read_excel(file_content, sheet_name=None, engine='openpyxl')
        text_parts = ["[EXCEL Data - COMPLETE FILE]"]
        
        for sheet_name, df in df_dict.items():
            # Skip empty sheets
            if df.empty:
                continue
                
            sheet_text = [f"[Sheet: {sheet_name}]"]
            
            headers = " | ".join([str(col) for col in df.columns if str(col)])
            sheet_text.append(f"Headers: {headers}")
            
            # Index ALL rows
            for idx, row in df.iterrows():
                row_text = " | ".join([str(val) for val in row.values if pd.notna(val) and str(val).strip()])
                if row_text:
                    sheet_text.append(row_text)
            
            text_parts.append("\n".join(sheet_text))
        
        return "\n\n".join(text_parts)
    except Exception as e:
        print(f"  Error extracting Excel: {e}")
        return ""


def extract_text_from_csv(file_content):
    """
    Extract text from CSV - NEW APPROACH: Store complete CSV data as single unit
    Instead of chunking for embeddings, we create a comprehensive summary + full data reference
    """
    try:
        file_content.seek(0)
        df = pd.read_csv(file_content, nrows=50000)  # Safety limit
        
        # NEW APPROACH: Create single comprehensive representation
        # This will be ONE chunk that references the complete data
        
        headers = " | ".join([str(col) for col in df.columns if str(col)])
        total_rows = len(df)
        
        # Create a rich summary that includes:
        # 1. File structure (headers, row count)
        # 2. Column names for semantic search
        # 3. Sample data (first/last rows)
        # 4. Statistical summary
        
        summary_parts = []
        
        # Header with metadata
        summary_parts.append(f"[CSV Data - COMPLETE FILE]")
        summary_parts.append(f"Headers: {headers}\n")
        summary_parts.append(f"Total Rows: {total_rows}\n")
        
        # Add column names repeatedly for better semantic search matching
        summary_parts.append(f"Column Names: {', '.join([str(col) for col in df.columns])}\n")
        
        # Check last few rows for totals/summary rows (common in financial CSVs)
        summary_parts.append("\n========== KEY TOTALS FROM FILE ==========")
        totals_found = False
        
        # Look at last 10 rows for total/summary rows
        last_rows = df.tail(10)
        for idx, row in last_rows.iterrows():
            # Check if this row contains total indicators
            first_col_value = str(row.iloc[0]).lower() if len(row) > 0 else ""
            if any(indicator in first_col_value for indicator in ['total', 'sum', 'grand', 'report total', '---', '===', 'summary']):
                totals_found = True
                # Format this row nicely
                row_parts = []
                for col_name, value in zip(df.columns, row.values):
                    if pd.notna(value) and str(value).strip():
                        # Format numeric values as currency if appropriate
                        try:
                            num_val = float(str(value).replace('$', '').replace(',', ''))
                            if num_val != 0:
                                col_str = str(col_name)
                                # Check if column looks like a month or currency column
                                if any(keyword in col_str.lower() for keyword in ['revenue', 'sales', 'amount', 'price', 'total', '-20', 'jan', 'feb', 'mar', 'apr', 'may', 'jun', 'jul', 'aug', 'sep', 'oct', 'nov', 'dec']):
                                    row_parts.append(f"  {col_name}: ${num_val:,.2f}")
                                else:
                                    row_parts.append(f"  {col_name}: {value}")
                        except:
                            row_parts.append(f"  {col_name}: {value}")
                
                if row_parts:
                    summary_parts.append(f"\n{first_col_value.upper()}:")
                    summary_parts.extend(row_parts)
        
        if not totals_found:
            # If no explicit total row, calculate column sums
            summary_parts.append("\n(Calculated column totals):")
            for col in df.columns:
                try:
                    numeric_col = pd.to_numeric(df[col], errors='coerce')
                    non_null_count = numeric_col.notna().sum()
                    if non_null_count > 0:
                        total = numeric_col.sum()
                        if total != 0:
                            col_str = str(col)
                            if any(keyword in col_str.lower() for keyword in ['revenue', 'sales', 'amount', 'price', 'cost', 'total', 'value', '-20', 'jan', 'feb', 'mar', 'apr', 'may', 'jun', 'jul', 'aug', 'sep', 'oct', 'nov', 'dec']) or '$' in col_str:
                                summary_parts.append(f"  {col}: ${total:,.2f}")
                except:
                    pass
        
        summary_parts.append("=" * 50 + "\n")
        
        # Add ALL the data (not chunked)
        # Format each row for readability
        summary_parts.append("\nCOMPLETE DATA:\n")
        for idx, row in df.iterrows():
            row_text = " | ".join([str(val) for val in row.values if pd.notna(val)])
            if row_text:
                summary_parts.append(row_text)
        
        # Add summary footer
        summary_parts.append(f"\n-------Report Total------- (All {total_rows} rows included above)")
        
        # Return complete CSV as single text block
        return "\n".join(summary_parts)
        
    except Exception as e:
        print(f"  Pandas CSV read failed ({e}), falling back to raw text.")
        file_content.seek(0)
        return file_content.read().decode('utf-8', errors='ignore')


def process_extracted_text(text: str, filename: str = "", source_type: str = "") -> Optional[str]:
    """
    Process extracted text through quality filtering and clarification
    Returns None if text quality is too low, otherwise returns processed text
    """
    if not text or len(text.strip()) < 10:
        return None
    
    # Get quality filter and clarification service
    quality_filter = get_quality_filter()
    clarification_service = get_clarification_service()
    
    # First, assess text quality
    quality_assessment = quality_filter.assess_text_quality(text, filename, source_type)
    
    if not quality_assessment['should_include']:
        logger.info(f"Text from {filename} excluded: {quality_assessment['reason']}")
        return None
    
    logger.info(f"Text from {filename} passed quality check (score: {quality_assessment['quality_score']:.2f})")
    
    # If quality is acceptable, try to clarify if needed
    processed_text = text
    if (clarification_service and clarification_service.model and 
        clarification_service.should_clarify_text(text, source_type)):
        
        try:
            clarified_text = clarification_service.clarify_text(
                text=text,
                filename=filename,
                source_type=source_type,
                context="Extracted document text"
            )
            
            # Re-assess quality of clarified text
            clarified_assessment = quality_filter.assess_text_quality(clarified_text, filename, f"Clarified {source_type}")
            
            if clarified_assessment['should_include'] and clarified_assessment['quality_score'] > quality_assessment['quality_score']:
                processed_text = clarified_text
                logger.info(f"Text from {filename} clarified and improved (quality: {quality_assessment['quality_score']:.2f} → {clarified_assessment['quality_score']:.2f})")
            else:
                logger.info(f"Clarification did not improve quality for {filename}, using original")
                
        except Exception as e:
            logger.warning(f"Failed to clarify text from {filename}: {e}")
    
    return processed_text


def extract_text_from_image(file_content, ocr_service, filename=""):
    """Extract text from image using OCR with quality filtering and AI clarification"""
    try:
        if not ocr_service:
            logger.warning("OCR service not available for image processing")
            return f"[IMAGE FILE: {filename}] - OCR not available"
        
        # Read image bytes
        if hasattr(file_content, 'read'):
            image_data = file_content.read()
        else:
            image_data = file_content
        
        # Use OCR to extract text
        ocr_result = ocr_service.extract_text(image_data, languages=OCR_LANGUAGES)
        
        if ocr_result.text.strip():
            # Format basic metadata
            metadata_parts = [f"[IMAGE FILE: {filename}]"]
            metadata_parts.append(f"OCR Confidence: {ocr_result.confidence:.2f}")
            if ocr_result.language:
                metadata_parts.append(f"Detected Language: {ocr_result.language}")
            
            # Process through quality filter and clarification
            processed_text = process_extracted_text(
                text=ocr_result.text,
                filename=filename,
                source_type="OCR Image"
            )
            
            if processed_text is None:
                logger.info(f"OCR text from {filename} filtered out due to low quality")
                return None  # Signal that this content should be skipped
            
            # Combine metadata with processed text
            metadata_header = "\n".join(metadata_parts)
            if processed_text != ocr_result.text:
                # Text was clarified
                return f"{metadata_header}\n\n--- AI-Enhanced Text ---\n{processed_text}"
            else:
                # Text was not clarified but passed quality check
                return f"{metadata_header}\n\n--- Extracted Text ---\n{processed_text}"
        else:
            logger.info(f"No text detected in image {filename} or confidence too low")
            return None  # Signal that this content should be skipped
            
    except Exception as e:
        logger.error(f"Error extracting text from image {filename}: {e}")
        return None  # Signal that this content should be skipped


def extract_text(file_content, mime_type, filename="", ocr_service=None):
    """Extract text from any supported file type including images"""
    try:
        # Check if it's an image file first
        if is_image_file(mime_type, filename):
            return extract_text_from_image(file_content, ocr_service, filename)
        
        # Existing document processing logic
        if mime_type == 'application/pdf':
            # First try to extract text normally from PDF
            pdf_reader = PdfReader(file_content)
            text = ""
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            
            # If no text was extracted (scanned PDF), try OCR
            if not text.strip() and ocr_service:
                logger.info(f"PDF '{filename}' appears to be scanned/image-based, attempting OCR...")
                try:
                    # Check if using Document AI (can process entire PDF directly)
                    if OCR_BACKEND == "documentai" and hasattr(ocr_service, 'process_pdf'):
                        logger.info(f"Using Document AI to process entire PDF: {filename}")
                        file_content.seek(0)
                        ocr_result = ocr_service.process_pdf(file_content.read())
                        combined_ocr_text = ocr_result.text
                    else:
                        # Fallback: Convert PDF pages to images and OCR each page
                        import fitz  # PyMuPDF for PDF to image conversion
                        file_content.seek(0)
                        pdf_document = fitz.open(stream=file_content.read(), filetype="pdf")
                        
                        ocr_text_parts = []
                        for page_num in range(len(pdf_document)):
                            # Convert page to image
                            page = pdf_document.load_page(page_num)
                            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for better OCR
                            image_data = pix.tobytes("png")
                            
                            # OCR the page image
                            ocr_result = ocr_service.extract_text(image_data, languages=OCR_LANGUAGES)
                            if ocr_result.text.strip():
                                ocr_text_parts.append(f"[Page {page_num + 1}]\n{ocr_result.text}")
                        
                        pdf_document.close()
                        combined_ocr_text = "\n\n".join(ocr_text_parts) if ocr_text_parts else ""
                    
                    if combined_ocr_text.strip():
                        
                        # Process through quality filter and clarification
                        processed_text = process_extracted_text(
                            text=combined_ocr_text,
                            filename=filename,
                            source_type="Scanned PDF OCR"
                        )
                        
                        if processed_text is None:
                            logger.info(f"OCR text from PDF {filename} filtered out due to low quality")
                            text = f"[SCANNED PDF: {filename}] - OCR text quality too low for indexing"
                        else:
                            if processed_text != combined_ocr_text:
                                text = f"[SCANNED PDF: {filename}] - Processed with OCR and AI Enhancement\n\n{processed_text}"
                                logger.info(f"Successfully enhanced PDF OCR text for {filename}")
                            else:
                                text = f"[SCANNED PDF: {filename}] - Processed with OCR\n\n{processed_text}"
                        
                        logger.info(f"OCR processed {len(ocr_text_parts)} pages of '{filename}'")
                    else:
                        text = f"[SCANNED PDF: {filename}] - No text could be extracted via OCR"
                        
                except ImportError:
                    logger.warning("PyMuPDF not installed - cannot OCR scanned PDFs. Install with: pip install PyMuPDF")
                    text = f"[SCANNED PDF: {filename}] - No text found and PyMuPDF not available for OCR"
                except Exception as e:
                    logger.error(f"Error during PDF OCR processing: {e}")
                    text = f"[SCANNED PDF: {filename}] - OCR processing failed: {str(e)}"
            
            return text
        
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            doc = DocxDocument(file_content)
            processed_paragraphs = []
            current_heading = ""
            
            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                
                # Check for built-in heading styles
                if para.style and para.style.name.lower().startswith('heading'):
                    current_heading = text
                    processed_paragraphs.append(f"\n\n[TOPIC: {current_heading}]\n")
                # Fallback heuristic
                elif para.style and para.style.name.lower().startswith('title') or (text.isupper() and len(text.split()) < 10):
                    current_heading = text
                    processed_paragraphs.append(f"\n\n[TOPIC: {current_heading}]\n")
                else:
                    # Prepend context
                    if current_heading:
                        processed_paragraphs.append(f"[{current_heading}] {text}")
                    else:
                        processed_paragraphs.append(text)
                        
            return "\n".join(processed_paragraphs)
        
        elif mime_type in ['application/vnd.openxmlformats-officedocument.presentationml.presentation',
                           'application/vnd.ms-powerpoint']:
            return extract_text_from_pptx(file_content)
        
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return extract_text_from_xlsx(file_content)
        
        elif mime_type == 'text/csv':
            return extract_text_from_csv(file_content)
        
        elif mime_type == 'text/plain':
            return file_content.read().decode('utf-8', errors='ignore')
        
        else:
            return ""
    except Exception as e:
        print(f"  Error extracting text: {e}")
        return ""


# --- ENHANCED: Semantic-aware text splitters ---

# Primary splitter - for child chunks (used in retrieval)
text_splitter = RecursiveCharacterTextSplitter(
    chunk_size=CHUNK_SIZE,
    chunk_overlap=CHUNK_OVERLAP,
    length_function=lambda text: len(text.split()),  # Use word count
    separators=["\n\n\n", "\n\n", "\n", ". ", "! ", "? ", "; ", ", ", " ", ""],  # Semantic boundaries
    is_separator_regex=False
)

# Parent splitter - for parent document retrieval (larger context)
parent_splitter = None
if USE_PARENT_DOCUMENT_RETRIEVAL:
    parent_splitter = RecursiveCharacterTextSplitter(
        chunk_size=PARENT_CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP * 2,  # More overlap for parent chunks
        length_function=lambda text: len(text.split()),
        separators=["\n\n\n", "\n\n", "\n", ". ", " ", ""],
        is_separator_regex=False
    )


def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP, return_parents=False):
    """
    Recursively split text into overlapping chunks using Langchain.
    CSV files are pre-chunked and split on chunk boundaries.
    
    Args:
        text: Text to chunk
        chunk_size: Target chunk size in words
        overlap: Overlap size in words
        return_parents: If True and parent retrieval enabled, return (child_chunks, parent_chunks)
    
    Returns:
        List of text chunks, or tuple of (child_chunks, parent_chunks) if return_parents=True
    """
    if not text or len(text.strip()) == 0:
        return [] if not return_parents else ([], [])
    
    # Check if this is a CSV with pre-defined chunk boundaries
    if "--- CSV CHUNK BOUNDARY ---" in text:
        # Split on boundaries and return as-is
        chunks = [chunk.strip() for chunk in text.split("--- CSV CHUNK BOUNDARY ---") if chunk.strip()]
        if return_parents and USE_PARENT_DOCUMENT_RETRIEVAL:
            return (chunks, chunks)  # Parent and child are the same for CSVs
        return chunks
    
    # Check if this is a complete file that shouldn't be chunked (Excel files)
    if "COMPLETE FILE" in text[:200]:
        # Return as single chunk without splitting
        if return_parents and USE_PARENT_DOCUMENT_RETRIEVAL:
            return ([text], [text])  # Parent and child are the same
        return [text]
    
    # Get child chunks for regular documents
    child_chunks = text_splitter.split_text(text)
    
    # Optionally get parent chunks
    if return_parents and USE_PARENT_DOCUMENT_RETRIEVAL and parent_splitter:
        parent_chunks = parent_splitter.split_text(text)
        return (child_chunks, parent_chunks)
    
    return child_chunks

# --- DELETED: Old _split_text_with_overlap and old chunk_text ---

